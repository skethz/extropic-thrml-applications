#!/usr/bin/env python3
"""Build a brief 2-slide deck: CUDA kernel vs thrml + trajectory-correction figure."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ---------------- Slide 1: comparison ----------------
s = prs.slides.add_slide(blank)

def textbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame

tf = textbox(s, 0.5, 0.25, 12.3, 1.0)
tf.text = "CUDA Ising Kernel  vs  Extropic thrml"
tf.paragraphs[0].font.size = Pt(32); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = NAVY
p = tf.add_paragraph()
p.text = "CR3BP L1 trajectory-correction maneuver solved as an Ising / QUBO problem  (N = 256 spins, K=16 maneuvers × 8 bits)"
p.font.size = Pt(14); p.font.color.rgb = GREY

rows = [
    ("", "CUDA Ising kernel", "Extropic thrml"),
    ("Implementation", "Hand-written CUDA (.cu), GH200", "Pure JAX — runs on CPU & GPU, unchanged"),
    ("Spin update", "Greedy single-spin flip (T = 0)", "SpinGibbsConditional heat-bath, β annealed"),
    ("Architecture", "num_reads chains × iters single-site steps", "Identical loop & parameters (drop-in)"),
    ("p_attempt (10 runs)", "1.0000", "1.0000"),
    ("best J / median J", "4.2e-7 / 4.6e-7", "5.1e-7 / 6.1e-7"),
    ("TTS@0.99", "1.6 s", "GPU 5.1 s   ·   CPU 182 s"),
    ("Reproducibility", "—", "Bit-identical CPU ↔ GPU (same seed)"),
]
nr, nc = len(rows), 3
gt = s.shapes.add_table(nr, nc, Inches(0.5), Inches(1.55), Inches(12.3), Inches(3.5)).table
gt.columns[0].width = Inches(3.0); gt.columns[1].width = Inches(4.4); gt.columns[2].width = Inches(4.9)
for r in range(nr):
    for c in range(nc):
        cell = gt.cell(r, c)
        cell.text = rows[r][c]
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(12.5)
        if r == 0:
            para.font.bold = True; para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            para.font.bold = (c == 0)
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r % 2 else RGBColor(0xFF, 0xFF, 0xFF)
            if c == 2:
                para.font.color.rgb = GREEN

tf = textbox(s, 0.5, 5.25, 12.3, 1.9)
tf.word_wrap = True
bullets = [
    "Same QUBO, same harness, same parameters — only the per-spin update is swapped: thrml is a drop-in replacement for the kernel.",
    "Full parity on the benchmark's success criterion (p = 1.0) and on solution quality (best/median objective in the same regime).",
    "A programmable thermodynamic sampler (β schedule) replaces the bespoke greedy kernel; thrml on GPU is ~36× faster than CPU and within ~3× of the hand-written CUDA kernel.",
]
for i, b in enumerate(bullets):
    par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    par.text = "•  " + b
    par.font.size = Pt(13); par.font.color.rgb = NAVY; par.space_after = Pt(6)

# ---------------- Slide 2: trajectory ----------------
s2 = prs.slides.add_slide(blank)
tf = textbox(s2, 0.5, 0.25, 12.3, 0.8)
tf.text = "Trajectory correction (planar CR3BP, Earth–Moon L1)"
tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = NAVY

s2.shapes.add_picture("traj_correction.png", Inches(0.55), Inches(1.25), width=Inches(12.2))

tf = textbox(s2, 0.5, 6.55, 12.3, 0.85)
tf.word_wrap = True
tf.text = ("A 0.38 km initial perturbation grows to a 573 km terminal error on the unstable L1 orbit. "
           "The thrml-solved maneuvers (16 impulsive Δv, total 26 m/s) cut the terminal error to ~51 km — "
           "an ~11× reduction. The objective targets the terminal state, so the corrected path deviates "
           "mid-flight, then converges.")
tf.paragraphs[0].font.size = Pt(12); tf.paragraphs[0].font.color.rgb = GREY

prs.save("cuda_vs_thrml.pptx")
print("wrote cuda_vs_thrml.pptx")
